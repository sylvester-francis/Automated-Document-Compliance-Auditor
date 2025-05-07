# Enhanced app/services/extraction_service.py

import io
import os
import re
import logging
import csv
from typing import Dict, List, Any, Optional, Union
from datetime import datetime

# Import existing dependencies
import PyPDF2
from PyPDF2 import PdfReader

logger = logging.getLogger(__name__)

class ExtractionService:
    """Service for extracting text from different document formats"""
    
    def __init__(self):
        """Initialize the extraction service"""
        # Load any necessary resources once at initialization
        self.supported_formats = {
            '.pdf': self.extract_from_pdf,
            '.docx': self.extract_from_docx,
            '.txt': self.extract_from_txt,
            '.html': self.extract_from_html,
            '.htm': self.extract_from_html,
            '.rtf': self.extract_from_rtf,
            '.csv': self.extract_from_csv,
            '.xlsx': self.extract_from_xlsx,
            '.xls': self.extract_from_xlsx,
            '.pptx': self.extract_from_pptx,
            '.ppt': self.extract_from_pptx,
            '.eml': self.extract_from_eml,
            '.json': self.extract_from_json,
            '.xml': self.extract_from_xml
        }
    
    def extract_text(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """
        Extract text from a document file
        
        Args:
            file_path: Path to the document file
            file_content: Optional bytes content of the file
            
        Returns:
            Dictionary with extracted text and metadata
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        result = {
            "text": "",
            "metadata": {},
            "format": ext[1:] if ext.startswith('.') else ext,
            "filename": os.path.basename(file_path),
            "extraction_time": datetime.now()
        }
        
        try:
            # Check if we have a handler for this format
            if ext in self.supported_formats:
                handler = self.supported_formats[ext]
                extracted_data = handler(file_path, file_content)
                result.update(extracted_data)
            else:
                # Try to read as text for unsupported formats
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                    result["text"] = text
                    result["metadata"]["warning"] = f"Unsupported file format: {ext}. Treating as plain text."
                except:
                    result["error"] = f"Unsupported file format: {ext}"
            
            # Calculate statistics about the extracted text
            if "text" in result and result["text"]:
                result["statistics"] = self._calculate_text_stats(result["text"])
                
            return result
        
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            result["error"] = str(e)
            return result
    
    def _calculate_text_stats(self, text: str) -> Dict[str, Any]:
        """Calculate statistics about the extracted text"""
        stats = {
            "char_count": len(text),
            "word_count": len(text.split()),
            "line_count": len(text.splitlines())
        }
        return stats
    
    def extract_from_pdf(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text and metadata from PDF file"""
        try:
            # Prepare file object
            if file_content:
                file_object = io.BytesIO(file_content)
            else:
                file_object = open(file_path, 'rb')
            
            # Try using PyMuPDF if available (better text extraction)
            try:
                import fitz  # PyMuPDF
                
                doc = fitz.open(stream=file_object.read(), filetype="pdf")
                
                # Extract text from each page
                text = ""
                page_texts = []
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    page_text = page.get_text()
                    text += page_text + "\n\n"
                    page_texts.append(page_text)
                
                # Extract metadata
                metadata = {
                    "page_count": len(doc),
                    "author": doc.metadata.get("author", ""),
                    "title": doc.metadata.get("title", ""),
                    "subject": doc.metadata.get("subject", ""),
                    "creator": doc.metadata.get("creator", ""),
                    "producer": doc.metadata.get("producer", ""),
                    "creation_date": doc.metadata.get("creationDate", ""),
                    "modification_date": doc.metadata.get("modDate", "")
                }
                
                # Check for form fields (potential medical forms)
                has_form = False
                for page in doc:
                    if page.widgets:
                        has_form = True
                        break
                metadata["has_form"] = has_form
                
                doc.close()
                
                # Add page-level data
                page_data = [{"page_num": i+1, "text_length": len(page_text)} 
                            for i, page_text in enumerate(page_texts)]
                
            except ImportError:
                # Fallback to PyPDF2
                file_object.seek(0)  # Reset file pointer
                
                pdf_reader = PdfReader(file_object)
                
                text = ""
                page_texts = []
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    text += page_text + "\n\n"
                    page_texts.append(page_text)
                
                # Extract metadata
                metadata = {}
                if pdf_reader.metadata:
                    for key, value in pdf_reader.metadata.items():
                        clean_key = key[1:] if key.startswith('/') else key
                        metadata[clean_key] = value
                metadata['page_count'] = len(pdf_reader.pages)
                
                # Add page-level data
                page_data = [{"page_num": i+1, "text_length": len(page_text)} 
                            for i, page_text in enumerate(page_texts)]
            
            # Close file if opened from path
            if not file_content:
                file_object.close()
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": metadata,
                "pages": page_data
            }
        
        except Exception as e:
            logger.error(f"Error extracting from PDF: {str(e)}")
            raise
    
    def extract_from_docx(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text and metadata from DOCX file"""
        try:
            import docx
            
            # Prepare file object
            if file_content:
                file_object = io.BytesIO(file_content)
                doc = docx.Document(file_object)
            else:
                doc = docx.Document(file_path)
            
            # Extract text from paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Extract text from tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            # Combined text
            text = "\n".join(paragraphs)
            
            # Add table content to text
            for table in tables:
                for row in table:
                    text += "\n" + " | ".join(row)
                text += "\n"
            
            # Extract metadata
            metadata = {
                "paragraph_count": len(paragraphs),
                "table_count": len(tables)
            }
            
            if hasattr(doc, 'core_properties'):
                props = doc.core_properties
                if hasattr(props, 'author'):
                    metadata["author"] = props.author
                if hasattr(props, 'title'):
                    metadata["title"] = props.title
                if hasattr(props, 'created'):
                    metadata["created"] = props.created
                if hasattr(props, 'modified'):
                    metadata["modified"] = props.modified
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": metadata,
                "paragraphs": paragraphs,
                "tables": tables
            }
        
        except ImportError:
            logger.warning("python-docx not installed, falling back to simple extraction")
            return {
                "text": "DOCX file cannot be processed (python-docx library not installed)",
                "metadata": {"error": "Missing dependency: python-docx"}
            }
        except Exception as e:
            logger.error(f"Error extracting from DOCX: {str(e)}")
            raise
    
    def extract_from_txt(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from TXT file"""
        try:
            if file_content:
                text = file_content.decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
            
            # Basic statistics
            lines = text.splitlines()
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": {
                    "format": "txt",
                    "size": len(text),
                    "line_count": len(lines)
                }
            }
        
        except Exception as e:
            logger.error(f"Error extracting from TXT: {str(e)}")
            raise
    
    def extract_from_html(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            
            # Read content
            if file_content:
                content = file_content.decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            # Parse HTML
            soup = BeautifulSoup(content, 'html.parser')
            
            # Extract title
            title = soup.title.string if soup.title else ""
            
            # Remove script and style elements
            for tag in soup(["script", "style", "noscript", "iframe", "head"]):
                tag.extract()
            
            # Get text
            text = soup.get_text(separator='\n')
            
            # Clean up text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            # Extract headers
            headers = []
            for i in range(1, 7):
                for header in soup.find_all(f'h{i}'):
                    headers.append({"level": i, "text": header.get_text().strip()})
            
            # Extract links
            links = []
            for link in soup.find_all('a', href=True):
                links.append({
                    "text": link.get_text().strip(),
                    "href": link['href']
                })
            
            # Extract metadata
            metadata = {
                "title": title,
                "header_count": len(headers),
                "link_count": len(links),
                "format": "html"
            }
            
            # Add meta tags
            meta_tags = {}
            for meta in soup.find_all('meta'):
                name = meta.get('name') or meta.get('property')
                if name:
                    meta_tags[name] = meta.get('content', '')
            
            if meta_tags:
                metadata["meta_tags"] = meta_tags
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": metadata,
                "headers": headers,
                "links": links
            }
        
        except ImportError:
            # BeautifulSoup not available
            if file_content:
                content = file_content.decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            # Simple regex-based extraction
            text = re.sub(r'<[^>]+>', ' ', content)
            text = re.sub(r'\s+', ' ', text).strip()
            
            return {
                "text": text,
                "metadata": {
                    "format": "html",
                    "warning": "Limited extraction (BeautifulSoup not installed)"
                }
            }
        except Exception as e:
            logger.error(f"Error extracting from HTML: {str(e)}")
            raise
    
    def extract_from_rtf(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from RTF file"""
        try:
            # Try using the striprtf library
            from striprtf.striprtf import rtf_to_text
            
            if file_content:
                content = file_content.decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            text = rtf_to_text(content)
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": {
                    "format": "rtf"
                }
            }
        except ImportError:
            # Fallback to a more basic approach
            if file_content:
                content = file_content.decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            # Simple regex to remove RTF commands
            text = re.sub(r'\\[a-z0-9]+', ' ', content)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\\\n', '\n', text)
            text = re.sub(r'\s+', ' ', text).strip()
            
            return {
                "text": text,
                "metadata": {
                    "format": "rtf",
                    "warning": "Basic extraction (striprtf library not installed)"
                }
            }
        except Exception as e:
            logger.error(f"Error extracting from RTF: {str(e)}")
            raise
    
    def extract_from_csv(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from CSV file"""
        try:
            rows = []
            
            if file_content:
                content = file_content.decode('utf-8', errors='ignore')
                rows = list(csv.reader(content.splitlines()))
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    rows = list(reader)
            
            # Convert CSV data to text
            if rows:
                headers = rows[0]
                
                # Format as text
                text = ""
                
                # Add headers
                text += ", ".join(headers) + "\n\n"
                
                # Add rows
                for row in rows[1:]:
                    row_text = ""
                    for i, cell in enumerate(row):
                        if i < len(headers):
                            row_text += f"{headers[i]}: {cell}, "
                    text += row_text.rstrip(", ") + "\n"
            else:
                text = "Empty CSV file"
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": {
                    "format": "csv",
                    "row_count": len(rows) - 1 if rows else 0,
                    "column_count": len(rows[0]) if rows else 0
                },
                "data": rows
            }
        except Exception as e:
            logger.error(f"Error extracting from CSV: {str(e)}")
            raise
    
    def extract_from_xlsx(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from Excel XLSX file"""
        try:
            import openpyxl
            
            # Load workbook
            if file_content:
                wb = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True, data_only=True)
            else:
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            
            # Extract text from worksheets
            text = ""
            sheet_data = {}
            
            for sheet in wb:
                sheet_text = f"Sheet: {sheet.title}\n\n"
                
                # Extract cell values
                rows = []
                for row in sheet.iter_rows():
                    row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                    rows.append(row_values)
                    sheet_text += " | ".join(row_values) + "\n"
                
                text += sheet_text + "\n\n"
                sheet_data[sheet.title] = rows
            
            # Basic metadata
            metadata = {
                "format": "xlsx",
                "sheets": list(wb.sheetnames),
                "sheet_count": len(wb.sheetnames)
            }
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": metadata,
                "sheet_data": sheet_data
            }
        except ImportError:
            # openpyxl not available
            return {
                "text": "Excel file cannot be processed (openpyxl library not installed)",
                "metadata": {
                    "format": "xlsx",
                    "error": "Missing dependency: openpyxl"
                }
            }
        except Exception as e:
            logger.error(f"Error processing Excel file: {str(e)}")
            raise
    
    def extract_from_pptx(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from PowerPoint PPTX file"""
        try:
            from pptx import Presentation
            
            # Load presentation
            if file_content:
                prs = Presentation(io.BytesIO(file_content))
            else:
                prs = Presentation(file_path)
            
            # Extract text from slides
            text = ""
            slides_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"Slide {i+1}:\n"
                slide_content = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                        slide_content.append(shape.text)
                
                text += slide_text + "\n"
                slides_data.append({
                    "number": i+1,
                    "content": slide_content
                })
            
            # Basic metadata
            metadata = {
                "format": "pptx",
                "slide_count": len(prs.slides)
            }
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": metadata,
                "slides": slides_data
            }
        except ImportError:
            # python-pptx not available
            return {
                "text": "PowerPoint file cannot be processed (python-pptx library not installed)",
                "metadata": {
                    "format": "pptx",
                    "error": "Missing dependency: python-pptx"
                }
            }
        except Exception as e:
            logger.error(f"Error processing PowerPoint file: {str(e)}")
            raise
    
    def extract_from_eml(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from email file (EML)"""
        try:
            import email
            from email import policy
            
            # Read content
            if file_content:
                msg = email.message_from_bytes(file_content, policy=policy.default)
            else:
                with open(file_path, 'rb') as f:
                    msg = email.message_from_binary_file(f, policy=policy.default)
            
            # Extract header info
            headers = {
                "From": str(msg.get("From", "")),
                "To": str(msg.get("To", "")),
                "Cc": str(msg.get("Cc", "")),
                "Subject": str(msg.get("Subject", "")),
                "Date": str(msg.get("Date", ""))
            }
            
            # Extract body text
            body = ""
            
            # Handle multipart messages
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    content_disposition = str(part.get("Content-Disposition", ""))
                    
                    # Skip attachments
                    if "attachment" in content_disposition:
                        continue
                    
                    # Get text content
                    if content_type == "text/plain":
                        try:
                            body += part.get_payload(decode=True).decode(errors='replace') + "\n"
                        except:
                            body += str(part.get_payload()) + "\n"
                    elif content_type == "text/html":
                        # Simple HTML to text conversion
                        try:
                            html = part.get_payload(decode=True).decode(errors='replace')
                            from bs4 import BeautifulSoup
                            soup = BeautifulSoup(html, 'html.parser')
                            body += soup.get_text(separator='\n') + "\n"
                        except ImportError:
                            # Simple regex-based extraction
                            html = part.get_payload(decode=True).decode(errors='replace')
                            body += re.sub(r'<[^>]+>', ' ', html) + "\n"
                        except:
                            body += str(part.get_payload()) + "\n"
            else:
                # Not multipart - just get the payload
                try:
                    body = msg.get_payload(decode=True).decode(errors='replace')
                except:
                    body = str(msg.get_payload())
            
            # Format as text
            text = f"From: {headers['From']}\n"
            text += f"To: {headers['To']}\n"
            if headers['Cc']:
                text += f"Cc: {headers['Cc']}\n"
            text += f"Subject: {headers['Subject']}\n"
            text += f"Date: {headers['Date']}\n\n"
            text += body
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": {
                    "format": "eml",
                    "headers": headers
                }
            }
        except Exception as e:
            logger.error(f"Error processing email file: {str(e)}")
            
            # Fallback - read as text
            if file_content:
                content = file_content.decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            return {
                "text": content,
                "metadata": {
                    "format": "eml",
                    "error": str(e)
                }
            }
    
    def extract_from_json(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from JSON file"""
        import json
        
        try:
            # Read content
            if file_content:
                data = json.loads(file_content.decode('utf-8', errors='ignore'))
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    data = json.load(f)
            
            # Convert JSON to text
            if isinstance(data, dict):
                # Format dictionary as text
                text = self._format_json_dict(data)
            elif isinstance(data, list):
                # Format list as text
                text = ""
                for i, item in enumerate(data):
                    if isinstance(item, dict):
                        text += f"Item {i+1}:\n{self._format_json_dict(item)}\n\n"
                    else:
                        text += f"Item {i+1}: {item}\n"
            else:
                # Just convert to string
                text = str(data)
            
            # Clean up text
            text = self._clean_text(text)
            
            return {
                "text": text,
                "metadata": {
                    "format": "json",
                    "type": type(data).__name__,
                    "data_size": len(json.dumps(data))
                },
                "data": data  # Include the parsed JSON data
            }
        except Exception as e:
            logger.error(f"Error parsing JSON file: {str(e)}")
            
            # Fallback to reading as text
            if file_content:
                content = file_content.decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            return {
                "text": content,
                "metadata": {
                    "format": "json",
                    "error": str(e)
                }
            }
    
    def _format_json_dict(self, data: Dict, prefix: str = "") -> str:
        """Helper to format JSON dictionary as readable text"""
        text = ""
        for key, value in data.items():
            if isinstance(value, dict):
                text += f"{prefix}{key}:\n{self._format_json_dict(value, prefix + '  ')}"
            elif isinstance(value, list):
                text += f"{prefix}{key}:\n"
                for i, item in enumerate(value):
                    if isinstance(item, dict):
                        text += f"{prefix}  Item {i+1}:\n{self._format_json_dict(item, prefix + '    ')}"
                    else:
                        text += f"{prefix}  - {item}\n"
            else:
                text += f"{prefix}{key}: {value}\n"
        return text
    
    def extract_from_xml(self, file_path: str, file_content: bytes = None) -> Dict[str, Any]:
        """Extract text from XML file"""
        try:
            import xml.etree.ElementTree as ET
            
            # Parse XML
            if file_content:
                tree = ET.ElementTree(ET.fromstring(file_content.decode('utf-8', errors='ignore')))
            else:
                tree = ET.parse(file_path)
            
            root = tree.getroot()
            
            # Extract text by recursively processing elements
            def extract_text_from_element(element, path=""):
                current_path = path + "/" + element.tag
                result = {
                    "path": current_path,
                    "text": element.text.strip() if element.text else "",
                    "attributes": element.attrib,
                    "children": []
                }
                
                # Add children
                for child in element:
                    result["children"].append(extract_text_from_element(child, current_path))
                
                return result
            
            # Build structured representation
            xml_structure = extract_text_from_element(root)
            
            # Convert to flat text
            def structure_to_text(structure, indent=0):
                text = " " * indent + structure["path"]
                if structure["attributes"]:
                    text += " " + str(structure["attributes"])
                if structure["text"]:
                    text += ": " + structure["text"]
                text += "\n"
                
                for child in structure["children"]:
                    text += structure_to_text(child, indent + 2)
                
                return text
            
            text = structure_to_text(xml_structure)
            
            # Clean up text
            text = self._clean_text(text)
            
            # Extract basic metadata
            metadata = {
                "root_tag": root.tag,
                "format": "xml",
                "element_count": len(root.findall(".//*")) + 1  # +1 for root
            }
            
            return {
                "text": text,
                "metadata": metadata,
                "structure": xml_structure
            }
        except Exception as e:
            logger.error(f"Error parsing XML file: {str(e)}")
            
            # Fallback to reading as text
            if file_content:
                content = file_content.decode('utf-8', errors='ignore')
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            return {
                "text": content,
                "metadata": {
                    "format": "xml",
                    "error": str(e)
                }
            }
    
    def _clean_text(self, text: str) -> str:
        """
        Clean text with common processing steps
        
        Args:
            text: Raw text to clean
            
        Returns:
            Cleaned text
        """
        if not text:
            return ""
        
        # Replace multiple newlines with a single one
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Replace multiple spaces with a single one
        text = re.sub(r' {2,}', ' ', text)
        
        # Remove null bytes
        text = text.replace('\x00', '')
        
        # Remove control characters (except newlines and tabs)
        text = re.sub(r'[\x01-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        
        return text.strip()

# Singleton instance
_extraction_service_instance = None

def get_extraction_service() -> ExtractionService:
    """
    Get or create the extraction service singleton instance
    
    Returns:
        ExtractionService instance
    """
    global _extraction_service_instance
    if _extraction_service_instance is None:
        _extraction_service_instance = ExtractionService()
    return _extraction_service_instance